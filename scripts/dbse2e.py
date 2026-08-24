#!/usr/bin/env python3
"""/dbse2e — end-to-end CITATION test for the Ask surface (#54).

Proves the thing the relevance floor (#53) is supposed to guarantee, on a realistic corpus:

  1. the CORRECT source IS cited for an on-topic question, and
  2. an INCORRECT/off-topic source does NOT pop up, and
  3. a SEMANTIC VARIATION ("is amazon bad?") still retrieves the relevant source whose
     polarity is the opposite ("Amazon ... is great").

How it works (deterministic, no Ollama needed): it builds a fresh in-memory QueryService
(HashingEmbedding + the default relevance floor — the SAME retrieval path the server uses),
ingests a few labelled fixture docs PLUS sampled rows from the 4 Kaggle datasets as realistic
distractors, then asserts on the citations `QueryService.answer()` returns — exactly what the
Ask UI shows under "Sources".

    python3 scripts/dbse2e.py            # run all cases
    python3 scripts/dbse2e.py -n 60      # 60 distractor rows per dataset (default 40)

Exit code is the number of failing cases (0 = all green), so CI can gate on it.
"""
from __future__ import annotations

import argparse
import csv
import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from dbsearch.adapters.local import (  # noqa: E402
    ExtractiveLlm, HashingEmbedding, InMemoryIdentity, InMemoryIndex,
    InMemoryObjectStore, InMemoryQueue, LocalRichExtractor,
)
from dbsearch.connectors.upload import UploadConnector  # noqa: E402
from dbsearch.pipeline.runner import run_ingestion  # noqa: E402
from dbsearch.query import QueryService  # noqa: E402

_KAGGLE = Path.home() / ".cache" / "kagglehub" / "datasets"

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _build_deck() -> bytes | None:
    """Build a small in-memory .pptx for the slide-citation proof. Returns None if python-pptx
    isn't installed (dbse2e stays runnable bare, like the Kaggle-dataset guard)."""
    try:
        from pptx import Presentation
    except Exception:
        return None
    prs = Presentation(); blank = prs.slide_layouts[6]
    slides = [
        "Company overview and mission",
        "Financial results revenue grew forty percent year over year",
        "Hiring plan and open roles",
    ]
    for t in slides:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(0, 0, 100, 100).text_frame.text = t
    import io as _io
    buf = _io.BytesIO(); prs.save(buf); return buf.getvalue()

# ---- dataset loaders: each yields (external_id, title, text) docs ---------------------------
# external_id is prefixed with the dataset key so must_cite / must_not_cite can match by source.

# Inline fallback distractors — used when the Kaggle datasets aren't on disk (they were archived
# out of ~/.cache; see .setaside_datasets/). Small but enough to keep all four cases meaningful:
# the amazon set includes the "vitality canned dog food" review T4 asserts on; the others are
# off-topic noise that must NOT be cited.
_INLINE = {
    "amazon-food": [
        ("amazon-food-1", "Amazon food review #1: Good Quality Dog Food",
         "Good Quality Dog Food. I have bought several of the Vitality canned dog food products and "
         "have found them all to be of good quality. The product looks more like a stew than a "
         "processed meat and my dog appreciates this product better than most."),
        ("amazon-food-2", "Amazon food review #2: Great coffee",
         "Great coffee. Rich bold flavour, arrives fresh, my favourite morning roast."),
        ("amazon-food-3", "Amazon food review #3: Tasty cat treats",
         "Tasty cat treats. My cat goes crazy for these, great value pack and good quality."),
    ],
    "airline": [
        ("airline-1", "Airline tweet (United, negative)", "@united my flight was delayed three hours and the service was bad."),
        ("airline-2", "Airline tweet (Delta, positive)", "@Delta great crew today, smooth boarding and on time, thank you."),
    ],
    "sarcasm": [
        ("sarcasm-1", "News headline #1", "area man passionate defender of what he imagines constitution to be"),
        ("sarcasm-2", "News headline #2", "local woman thrilled to discover meeting could have been an email"),
    ],
    "sms": [
        ("sms-1", "SMS message #1 (ham)", "Hey are we still on for dinner tonight at 7?"),
        ("sms-2", "SMS message #2 (spam)", "WINNER! You have won a free prize, text CLAIM to 80085 now."),
    ],
}


def _amazon(n: int):
    p = _KAGGLE / "snap/amazon-fine-food-reviews/versions/2/Reviews.csv"
    if not p.exists():
        yield from _INLINE["amazon-food"][:n]
        return
    with p.open(encoding="utf-8", errors="ignore", newline="") as f:
        for i, row in enumerate(csv.DictReader(f)):
            if i >= n:
                break
            txt = f"{row.get('Summary', '')}. {row.get('Text', '')}"
            yield f"amazon-food-{row.get('Id', i)}", f"Amazon food review #{row.get('Id', i)}: {row.get('Summary', '')[:60]}", txt


def _airline(n: int):
    p = _KAGGLE / "crowdflower/twitter-airline-sentiment/versions/4/Tweets.csv"
    if not p.exists():
        yield from _INLINE["airline"][:n]
        return
    with p.open(encoding="utf-8", errors="ignore", newline="") as f:
        for i, row in enumerate(csv.DictReader(f)):
            if i >= n:
                break
            yield f"airline-{row.get('tweet_id', i)}", f"Airline tweet ({row.get('airline', '')}, {row.get('airline_sentiment', '')})", row.get("text", "")


def _sarcasm(n: int):
    p = _KAGGLE / "rmisra/news-headlines-dataset-for-sarcasm-detection/versions/2/Sarcasm_Headlines_Dataset.json"
    if not p.exists():
        yield from _INLINE["sarcasm"][:n]
        return
    with p.open(encoding="utf-8", errors="ignore") as f:
        for i, line in enumerate(f):
            if i >= n:
                break
            line = line.strip()
            if not line:
                continue
            rec = json.loads(line)
            yield f"sarcasm-{i}", f"News headline #{i}", rec.get("headline", "")


def _sms(n: int):
    p = _KAGGLE / "uciml/sms-spam-collection-dataset/versions/1/spam.csv"
    if not p.exists():
        yield from _INLINE["sms"][:n]
        return
    with p.open(encoding="latin-1", newline="") as f:
        r = csv.reader(f)
        next(r, None)  # header v1,v2
        for i, row in enumerate(r):
            if i >= n or len(row) < 2:
                break
            yield f"sms-{i}", f"SMS message #{i} ({row[0]})", row[1]


LOADERS = {"amazon-food": _amazon, "airline": _airline, "sarcasm": _sarcasm, "sms": _sms}

# ---- fixtures: the docs whose citation behaviour we assert on -------------------------------
FIXTURES = [
    ("demo-handbook", "Staff Handbook",
     "Our holiday and expenses policy: all staff receive 25 days annual leave plus public "
     "holidays. Expenses are reimbursed monthly against receipts submitted in the portal. "
     "Remote work is permitted up to three days per week."),
    ("jay-shetty", "8 Rules Of Love by Jay Shetty.pdf",
     "8 Rules of Love by Jay Shetty. Jay R. Shetty teaches how to love, navigate relationships, "
     "heartbreak and breakups. He founded the Jay Shetty Certification School to train life coaches."),
    # sentiment fixture for the semantic-variation case: POSITIVE about Amazon (the show/service).
    ("amazon-prime-review", "Amazon Prime — viewer review",
     "Amazon is a good show service. Amazon Prime original series have excellent acting, "
     "compelling plots and superb production quality. I think Amazon is great and recommend it."),
]

# ---- test cases: query + which source ids MUST / MUST NOT appear in the citations ----------
# must_cite/must_not_cite match a citation if the substring is in its doc id (the source).
CASES = [
    {
        "name": "T1 correct citation: holiday/expenses -> Staff Handbook only",
        "query": "What is our holiday and expenses policy?",
        "must_cite": ["demo-handbook"],
        "must_not_cite": ["jay-shetty", "amazon-food", "airline", "sarcasm", "sms"],
    },
    {
        "name": "T2 correct source + incorrect absent: who is jay shetty",
        "query": "who is jay shetty",
        "must_cite": ["jay-shetty"],
        "must_not_cite": ["demo-handbook", "amazon-food", "airline", "sms"],
    },
    {
        "name": "T3 semantic variation: 'is amazon bad' retrieves the positive Amazon review",
        "query": "is amazon bad",
        "must_cite": ["amazon-prime-review"],
        "must_not_cite": ["demo-handbook", "jay-shetty"],
    },
    {
        "name": "T4 grounded-in-data: dog food quality -> an Amazon food review, not airline/sms",
        "query": "what do reviewers say about vitality canned dog food quality",
        "must_cite": ["amazon-food"],
        "must_not_cite": ["airline", "sms", "jay-shetty", "demo-handbook"],
    },
]


def build_service(n_per_dataset: int):
    store, queue = InMemoryObjectStore(), InMemoryQueue()
    index = InMemoryIndex(store)
    embedder = HashingEmbedding()
    extractor = LocalRichExtractor()
    docs = list(FIXTURES)
    counts = {"fixtures": len(FIXTURES)}
    for key, loader in LOADERS.items():
        loaded = list(loader(n_per_dataset))
        counts[key] = len(loaded)
        docs.extend(loaded)
    for ext, _title, text in docs:
        conn = UploadConnector("t", ext, ext, (text or "").encode(), "text/plain", ["all-staff"], "")
        run_ingestion(conn, queue, store, extractor, embedder, index)
    deck = _build_deck()
    deck_available = deck is not None
    if deck_available:
        conn = UploadConnector("t", "acme-deck", "Acme Deck.pptx", deck, PPTX_MIME, ["all-staff"], "upload://deck.pptx")
        run_ingestion(conn, queue, store, extractor, embedder, index)
    identity = InMemoryIdentity({"alice": ["all-staff"]})
    qs = QueryService(index, identity, embedder, ExtractiveLlm(), store)  # default floor=0.6
    return qs, counts, len(docs), deck_available


def run(n_per_dataset: int) -> int:
    qs, counts, total, deck_available = build_service(n_per_dataset)
    print(f"Corpus: {total} docs  ({counts})\n")
    fails = 0
    for c in CASES:
        res = qs.answer("alice", c["query"])
        cited = [cit["doc"] for cit in res.citations]
        missing = [m for m in c["must_cite"] if not any(m in d for d in cited)]
        leaked = [m for m in c["must_not_cite"] if any(m in d for d in cited)]
        ok = not missing and not leaked
        fails += 0 if ok else 1
        print(f"  {'PASS' if ok else 'FAIL'}  {c['name']}")
        print(f"        query    : {c['query']!r}")
        print(f"        cited    : {cited}")
        if missing:
            print(f"        MISSING  : {missing} (expected cited)")
        if leaked:
            print(f"        LEAKED   : {leaked} (should NOT be cited)")

    if deck_available:
        res = qs.answer("alice", "what were the financial results and revenue growth")
        deck_cites = [c for c in res.citations if c["doc"] == "acme-deck"]
        slide_ok = any((c.get("locator") or {}).get("kind") == "slide"
                       and (c.get("locator") or {}).get("n") == 2 for c in deck_cites)
        fails += 0 if slide_ok else 1
        print(f"  {'PASS' if slide_ok else 'FAIL'}  T5 slide-locator: revenue query cites Acme Deck slide 2")
        print(f"        citations: {[(c['doc'], c.get('locator')) for c in res.citations]}")
    else:
        print("  SKIP  T5 slide-locator: python-pptx not installed (pip install '.[server]')")

    print(f"\n{'ALL DBSE2E CASES PASSED' if not fails else f'{fails} CASE(S) FAILED'} "
          f"({len(CASES) - fails}/{len(CASES)} green)")
    return fails


def main() -> None:
    ap = argparse.ArgumentParser(description="/dbse2e — end-to-end citation test")
    ap.add_argument("-n", "--rows", type=int, default=40, help="distractor rows per dataset (default 40)")
    args = ap.parse_args()
    if not _KAGGLE.exists():
        print(f"WARNING: {_KAGGLE} not found — running with fixtures only (no dataset distractors).")
    sys.exit(run(args.rows))


if __name__ == "__main__":
    main()
