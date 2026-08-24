"""PPTX -> one segment per slide (shape text + speaker notes). python-pptx lazy-imported."""
from __future__ import annotations

import io

from dbsearch.core.models import Segment


def segment(data: bytes) -> list[Segment]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[Segment] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None and notes_tf.text.strip():
                parts.append(f"Notes: {notes_tf.text}")
        text = "\n".join(parts).strip()
        if text:
            out.append(Segment(text=text, locator={"kind": "slide", "n": i}))
    return out
