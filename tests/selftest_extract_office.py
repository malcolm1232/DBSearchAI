import io
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from dbsearch.adapters.local.extract import pptx as pptx_parser
from dbsearch.adapters.local.extract import docx as docx_parser


def _build_pptx() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title in ["Revenue outlook Q3", "Hiring plan headcount 40", "Risks and mitigations"]:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
        tb.text = title
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def _build_docx() -> bytes:
    from docx import Document
    d = Document()
    d.add_heading("Pricing", level=1)
    d.add_paragraph("Our day rate is 2000 per consultant.")
    d.add_heading("Timeline", level=1)
    d.add_paragraph("Delivery in twelve weeks.")
    buf = io.BytesIO(); d.save(buf); return buf.getvalue()


def test_pptx_segments_per_slide():
    segs = pptx_parser.segment(_build_pptx())
    assert len(segs) == 3
    assert segs[1].locator == {"kind": "slide", "n": 2}
    assert "Hiring plan" in segs[1].text


def test_docx_segments_per_section():
    segs = docx_parser.segment(_build_docx())
    assert [s.locator["kind"] for s in segs] == ["section", "section"]
    assert segs[0].locator["heading"] == "Pricing"
    assert "day rate" in segs[0].text
    assert segs[1].locator["n"] == 2


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("test_"):
            fn(); print("PASS", name)
