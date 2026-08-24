import io
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from dbsearch.adapters.local import (
    ExtractiveLlm, HashingEmbedding, InMemoryIdentity, InMemoryIndex,
    InMemoryObjectStore, InMemoryQueue, LocalRichExtractor,
)
from dbsearch.connectors.upload import UploadConnector
from dbsearch.pipeline.runner import run_ingestion
from dbsearch.query import QueryService

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _deck() -> bytes:
    from pptx import Presentation
    prs = Presentation(); blank = prs.slide_layouts[6]
    for t in ["Company overview intro", "Revenue outlook grew forty percent", "Team and hiring"]:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(0, 0, 100, 100).text_frame.text = t
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def _service():
    store, queue = InMemoryObjectStore(), InMemoryQueue()
    index = InMemoryIndex(store)
    embedder = HashingEmbedding()
    conn = UploadConnector("t", "deck", "Deck.pptx", _deck(), PPTX_MIME, ["all-staff"], "upload://Deck.pptx")
    run_ingestion(conn, queue, store, embedder=embedder, extractor=LocalRichExtractor(), index=index)
    identity = InMemoryIdentity({"alice": ["all-staff"]})
    return QueryService(index, identity, embedder, ExtractiveLlm(), store, tenant_id="t")


def test_citation_carries_slide_locator():
    qs = _service()
    res = qs.answer("alice", "what was the revenue outlook")
    locs = [c.get("locator", {}) for c in res.citations if c["doc"] == "deck"]
    assert any(l.get("kind") == "slide" and l.get("n") == 2 for l in locs), res.citations


if __name__ == "__main__":
    test_citation_carries_slide_locator(); print("PASS test_citation_carries_slide_locator")
