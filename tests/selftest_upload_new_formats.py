import io
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))
sys.path.insert(0, str(Path(__file__).resolve().parent))

import os
os.environ.setdefault("DBSEARCH_DEV_AUTH", "1")
from _upload import upload_and_wait
from fastapi.testclient import TestClient
from dbsearch.server.app import app

client = TestClient(app)


def _csv_bytes():
    return b"name,role\nAlice,Partner\nBob,Analyst\n"


def test_csv_upload_indexes_rows():
    files = {"file": ("team.csv", _csv_bytes(), "text/csv")}
    # #917: the real chunk count is read back via the LAW 2-trimmed segments endpoint,
    # so the uploader must be on the ACL to see their own count.
    out = upload_and_wait(client, {"X-DBSearch-User": "test_user"}, files=files,
                          data={"acl": ["all-staff", "test_user"], "title": "Team"})
    assert out["job"]["status"] == "succeeded", out["job"]
    assert out["chunk_count"] >= 1


def test_pptx_upload_indexes_slides():
    from pptx import Presentation
    prs = Presentation(); blank = prs.slide_layouts[6]
    prs.slides.add_slide(blank).shapes.add_textbox(0, 0, 100, 100).text_frame.text = "Quarterly revenue update"
    buf = io.BytesIO(); prs.save(buf)
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    files = {"file": ("deck.pptx", buf.getvalue(), mime)}
    out = upload_and_wait(client, {"X-DBSearch-User": "test_user"}, files=files,
                          data={"acl": ["all-staff", "test_user"]})
    assert out["job"]["status"] == "succeeded", out["job"]
    assert out["chunk_count"] >= 1


if __name__ == "__main__":
    test_csv_upload_indexes_rows(); print("PASS csv")
    test_pptx_upload_indexes_slides(); print("PASS pptx")
